import streamlit as st
import anthropic
import json
import base64
import os
from dotenv import load_dotenv
from io import BytesIO
import re

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_DATA_LABEL_POSITION

from PIL import Image

# ------------- 1) SETUP -------------
load_dotenv()
st.set_page_config(page_title="PEEL", layout="wide")

# Custom CSS for styling
st.markdown("""
    <style>
    .block-container {
        padding-top: 1rem;
        padding-bottom: 0rem;
    }
    div[data-testid="stFileUploader"] {
        width: 100%;
        padding: 1rem;
    }
    /* Custom styling for buttons */
    .stButton > button, .stDownloadButton > button {
        background-color: #e7fd7d !important;
        color: #544ff0 !important;
        font-size: 1.2em !important;
        padding: 0.8em 1.6em !important;
        transition: all 0.3s ease !important;
    }
    .stButton > button:hover, .stDownloadButton > button:hover {
        background-color: #d9fc5c !important;
        transform: scale(1.05);
    }
    /* Error message styling */
    div[data-baseweb="notification"] {
        margin-top: 50px;
    }
    </style>
""", unsafe_allow_html=True)

API_KEY = os.getenv("ANTHROPIC_API_KEY")
if not API_KEY:
    st.error("Anthropic API key not found. Please set it in .env or environment variables.")
    st.stop()

client = anthropic.Anthropic(
    api_key=API_KEY,
    default_headers={"anthropic-beta": "pdfs-2024-09-25"}
)
MODEL_NAME = "claude-3-5-sonnet-20241022"

EXETER_TEMPLATE_PATH = "exetertemplate2.pptx"
PROMPT_FILE = "prompt.txt"

def load_prompt_text(prompt_path):
    with open(prompt_path, "r", encoding="utf-8") as f:
        return f.read()

def get_slide_layouts(pptx_path):
    prs = Presentation(pptx_path)
    layout_info = []
    slide_master = prs.slide_masters[0]
    for layout in slide_master.slide_layouts:
        placeholders = []
        for shape in layout.placeholders:
            placeholders.append({
                "placeholder_idx": shape.placeholder_format.idx,
                "placeholder_name": shape.name,
                "shape_type": str(shape.shape_type)
            })
        layout_info.append({
            "layout_name": layout.name,
            "placeholders": placeholders
        })
    return layout_info

def build_flexible_prompt(layout_info, simplification_level, image_filenames):
    prompt_template = load_prompt_text(PROMPT_FILE)
    layout_info_json = json.dumps(layout_info, indent=2)
    image_filenames_json = json.dumps(image_filenames, indent=2)
    
    prompt_filled = (
        prompt_template
        .replace("{{SIMPLIFICATION_LEVEL}}", str(simplification_level))
        .replace("{{LAYOUT_INFO_JSON}}", layout_info_json)
        .replace("{{AVAILABLE_IMAGES}}", image_filenames_json)
    )
    return prompt_filled

def call_claude_for_slides(pdf_bytes, layout_info, simplification_level, image_filenames):
    final_prompt = build_flexible_prompt(layout_info, simplification_level, image_filenames)
    pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")
    
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": pdf_b64
                    }
                },
                {
                    "type": "text",
                    "text": final_prompt
                }
            ]
        }
    ]

    response = client.messages.create(
        model=MODEL_NAME,
        messages=messages,
        max_tokens=8192
    )

    assistant_reply = response.content[0].text
    json_pattern = r"<json_output>\s*(\{.*?\})\s*</json_output>"
    match = re.search(json_pattern, assistant_reply, re.DOTALL)
    
    if match:
        try:
            raw_json = match.group(1)
            return json.loads(raw_json)
        except json.JSONDecodeError:
            st.error("Failed to generate PowerPoint")
            return None
    else:
        try:
            return json.loads(assistant_reply)
        except:
            st.error("Failed to generate PowerPoint")
            return None

def find_placeholder_by_idx(slide, idx):
    """
    Locate the placeholder shape that has the given placeholder_format.idx
    """
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == idx:
            return shp
    return None

def add_picture_cover_cropped(slide, image_data, left, top, width, height):
    """
    Adds an image to the slide *covering* the specified bounding box
    without stretching. Crops any overflow to achieve a center-crop “cover”.
    """
    # Load image for ratio
    img = Image.open(image_data)
    w, h = img.size
    image_ratio = w / float(h)
    placeholder_ratio = width / float(height)

    # Create the Picture shape with the bounding box
    pic = slide.shapes.add_picture(image_data, left, top, width=width, height=height)
    
    # pic.crop_left, pic.crop_right, etc., are fractions of the original image
    if image_ratio > placeholder_ratio:
        # The image is relatively wider => fill height, crop left/right
        crop_amount = (1.0 - (placeholder_ratio / image_ratio)) / 2.0
        pic.crop_left = crop_amount
        pic.crop_right = crop_amount
        pic.crop_top = 0
        pic.crop_bottom = 0
    else:
        # The image is relatively taller => fill width, crop top/bottom
        crop_amount = (1.0 - (image_ratio / placeholder_ratio)) / 2.0
        pic.crop_top = crop_amount
        pic.crop_bottom = crop_amount
        pic.crop_left = 0
        pic.crop_right = 0

    return pic

def create_slides_from_json(prs, slides_json, layout_info, uploaded_images=None):
    """
    Converts the JSON specification into actual PowerPoint slides.
    - Text placeholders (with or without bullet points)
    - Charts (donut, comparison bars, trend line)
    - Images using a “cover + crop” approach
    """
    slide_master = prs.slide_masters[0]
    layout_name_map = {
        info["layout_name"]: layout_obj
        for info, layout_obj in zip(layout_info, slide_master.slide_layouts)
    }

    for slide_def in slides_json.get("slides", []):
        layout_name = slide_def.get("layout_name")
        if not layout_name:
            continue

        layout = layout_name_map.get(layout_name)
        if not layout:
            continue

        slide = prs.slides.add_slide(layout)
        placeholder_defs = slide_def.get("placeholders", {})

        for placeholder_idx, content in placeholder_defs.items():
            try:
                ph_idx = int(placeholder_idx)
            except ValueError:
                continue

            shape = find_placeholder_by_idx(slide, ph_idx)
            if not shape:
                continue

            # If it's a dict, check if it's chart, image, or text data
            if isinstance(content, dict):
                # --- Chart placeholders ---
                if "chart_type" in content and "chart_data" in content:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    
                    chart_type = content["chart_type"]
                    chart_data = content["chart_data"]
                    
                    if chart_type == "donut":
                        _add_donut_chart(slide, left, top, width, height, chart_data)
                    elif chart_type == "comparison_bars":
                        _add_comparison_bars_chart(slide, left, top, width, height, chart_data)
                    elif chart_type == "trend_line":
                        _add_trend_line_chart(slide, left, top, width, height, chart_data)
                    continue

                # --- Image placeholders (cover + crop) ---
                if "image_key" in content and uploaded_images:
                    img_key = content["image_key"]
                    if img_key in uploaded_images:
                        # Remove the original placeholder shape
                        sp = shape._element
                        sp.getparent().remove(sp)

                        # Add the image in “cover” mode, with center-crop
                        add_picture_cover_cropped(
                            slide,
                            BytesIO(uploaded_images[img_key]),
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )
                    continue

                # --- Text placeholders with possible bullet array ---
                text_val = content.get("text", "")
                bullet_vals = content.get("bullets", [])
                
                tf = shape.text_frame
                tf.word_wrap = True

                # Clear existing paragraphs
                for _ in range(len(tf.paragraphs[1:])):
                    p = tf.paragraphs[-1]._p
                    p.getparent().remove(p)

                if text_val and not bullet_vals:
                    # Single block of text, no bullets
                    tf.text = text_val
                    p = tf.paragraphs[0]
                    p.bullet = False
                    p.level = 0

                elif not text_val and len(bullet_vals) == 1:
                    # A single bullet, but treat it as plain text
                    tf.text = bullet_vals[0]
                    p = tf.paragraphs[0]
                    p.bullet = False
                    p.level = 0

                else:
                    # Possibly text + multiple bullets
                    if text_val:
                        tf.text = text_val
                        p = tf.paragraphs[0]
                        p.bullet = False
                        p.level = 0
                    else:
                        tf.text = ""

                    if len(bullet_vals) > 1:
                        for bullet_text in bullet_vals:
                            bp = tf.add_paragraph()
                            bp.text = bullet_text
                            # if truly want bullets, do bp.bullet = True,
                            # but your request is "still got bullet points" so let's keep them off:
                            bp.bullet = False
                            bp.level = 0

            # If content is just a simple string
            elif isinstance(content, str):
                tf = shape.text_frame
                tf.word_wrap = True
                tf.text = content

                p = tf.paragraphs[0]
                p.bullet = False
                p.level = 0

    return prs

def _add_donut_chart(slide, left, top, width, height, chart_data):
    title = chart_data.get("title", "Distribution")
    data_list = chart_data.get("data", [])
    
    if not data_list:
        data_list = [{"category": "No Data", "value": 100}]
    
    categories = []
    values = []
    for item in data_list:
        category = item["category"]
        value = item["value"]
        category_with_percent = f"{category} ({value}%)"
        categories.append(category_with_percent)
        values.append(value)

    cd = ChartData()
    cd.categories = categories
    cd.add_series("Distribution", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        left, top, width, height,
        cd
    ).chart

    chart.plots[0].donut_hole_size = 60
    chart.has_title = True
    chart.chart_title.text_frame.text = title

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0"%"'
    data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    
    return chart

def _add_comparison_bars_chart(slide, left, top, width, height, chart_data):
    title = chart_data.get("title", "Comparison")
    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])
    x_axis = chart_data.get("x_axis", "Categories")
    y_axis = chart_data.get("y_axis", "Values")
    
    formatted_labels = [label.replace(" ", "\n") for label in labels]
    
    cat_data = CategoryChartData()
    cat_data.categories = formatted_labels
    cat_data.add_series(y_axis, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        cat_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title

    category_axis = chart.category_axis
    value_axis = chart.value_axis
    
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = y_axis
    category_axis.has_title = True
    category_axis.axis_title.text_frame.text = x_axis

    plot = chart.plots[0]
    plot.has_data_labels = True
    
    for series in plot.series:
        series.has_data_labels = True
        for point in series.points:
            point.data_label.show_value = True
            point.data_label.number_format = '0"%"'
            point.data_label.font.bold = True
            point.data_label.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    return chart

def _add_trend_line_chart(slide, left, top, width, height, chart_data):
    title = chart_data.get("title", "Trend")
    dates = chart_data.get("dates", [])
    values = chart_data.get("values", [])
    x_axis = chart_data.get("x_axis", "Time")
    y_axis = chart_data.get("y_axis", "Value")
    
    formatted_dates = [date.replace(" ", "\n") for date in dates]
    
    cat_data = CategoryChartData()
    cat_data.categories = formatted_dates
    cat_data.add_series(y_axis, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        left, top, width, height,
        cat_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title

    category_axis = chart.category_axis
    value_axis = chart.value_axis
    
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = y_axis
    category_axis.has_title = True
    category_axis.axis_title.text_frame.text = x_axis

    plot = chart.plots[0]
    plot.has_data_labels = True
    
    for series in plot.series:
        series.has_data_labels = True
        series.smooth = True  # Make the line smooth
        for point in series.points:
            point.data_label.show_value = True
            point.data_label.number_format = '0.0'
            point.data_label.font.bold = True
            point.data_label.position = XL_DATA_LABEL_POSITION.ABOVE

    return chart

def main():
    # Sidebar authentication
    with st.sidebar:
        st.title("Authentication")
        password = st.text_input("Enter Password", type="password")
        
        # Updated disclaimer text
        st.markdown("---")
        st.markdown("""
        **Disclaimer:**
        
        This tool uses AI to summarise medical documents, but it may contain errors or omissions. 
        Always verify critical information with trusted medical sources and consult a qualified 
        professional before making any decisions based on the content generated. The AI does not 
        provide medical advice, diagnosis, or treatment. Use at your own discretion.
        """)

    # Check password
    if password != "P33L2025":
        st.markdown("<div style='padding-top: 50px'></div>", unsafe_allow_html=True)
        st.error("Please enter the correct password to access the application.")
        return

    # Add space before logo
    st.markdown("<div style='padding-top: 10px'></div>", unsafe_allow_html=True)
    
    # Add logo and top space
    col1, col2, col3 = st.columns([1, 4, 1])
    with col1:
        st.image("Peel-Lemon.svg", width=100)
    
    # Add space before first title
    st.markdown("<div style='padding-top: 20px'></div>", unsafe_allow_html=True)
    st.markdown("## 1. Create your presentation")

    if "layout_info" not in st.session_state:
        st.session_state.layout_info = get_slide_layouts(EXETER_TEMPLATE_PATH)

    st.header("Select simplification level")
    simplification_level = st.select_slider(
        "",
        options=list(range(1, 11)),
        format_func=lambda x: "Academic" if x == 1 else "Patient" if x == 10 else f"Level {x}",
        value=5
    )

    st.header("Upload pdf")
    uploaded_pdf = st.file_uploader("", type=["pdf"])

    st.header("Upload images")
    uploaded_images_files = st.file_uploader(
        "",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True
    )

    uploaded_images = {}
    image_filenames = []
    if uploaded_images_files:
        for img in uploaded_images_files:
            uploaded_images[img.name] = img.read()
            image_filenames.append(img.name)

    # Replace the current "Let's Peel" button code with this:

    # Create two columns for the button and spinner
    button_col, spinner_col = st.columns([3, 1])
    
    with button_col:
        generate_slides = st.button("Let's Peel", use_container_width=True)
    
    # Create a placeholder for the spinner in the second column
    spinner_placeholder = spinner_col.empty()
    
    if uploaded_pdf and generate_slides:
        # Show spinner in the placeholder
        with spinner_placeholder.container():
            with st.spinner(""):
                pdf_bytes = uploaded_pdf.read()
                result = call_claude_for_slides(
                    pdf_bytes,
                    st.session_state.layout_info,
                    simplification_level,
                    image_filenames
                )
                if result:
                    st.session_state.slides_json = result
                    st.success("PowerPoint generated!")
                else:
                    st.error("Failed to generate PowerPoint")
        
        # Add space after Let's Peel button
        st.markdown("<div style='padding-bottom: 50px'></div>", unsafe_allow_html=True)
    if "slides_json" in st.session_state and st.session_state.slides_json:
        st.markdown("## 2. Edit your slides")

        slides = st.session_state.slides_json.get("slides", [])

        # Create slide options list with titles
        slide_options = []
        for i, slide in enumerate(slides):
            # Find a title placeholder if it exists
            title = None
            for ph_idx, content in slide.get("placeholders", {}).items():
                placeholder_info = None
                for layout in st.session_state.layout_info:
                    if layout['layout_name'] == slide.get("layout_name", ""):
                        for ph in layout['placeholders']:
                            if str(ph['placeholder_idx']) == str(ph_idx):
                                placeholder_info = ph
                                break
                if placeholder_info and "title" in placeholder_info['placeholder_name'].lower():
                    if isinstance(content, dict):
                        title = content.get("text", "")
                    else:
                        title = content
                    break
            
            label = f"Slide {i+1}"
            if title:
                label += f" - {title}"
            slide_options.append(label)

        selected_slide_label = st.selectbox(
            "Choose a slide:",
            options=slide_options,
            index=0,
            key="selected_slide_dropdown"
        )
        selected_slide_index = slide_options.index(selected_slide_label)

        st.markdown(f"### Editing Slide {selected_slide_index + 1}")

        selected_slide = slides[selected_slide_index]
        placeholders = selected_slide.get("placeholders", {})

        # Let user edit text/bullet placeholders
        for ph_idx, content in placeholders.items():
            placeholder_info = None
            for layout in st.session_state.layout_info:
                if layout['layout_name'] == selected_slide.get("layout_name", ""):
                    for ph in layout['placeholders']:
                        if str(ph['placeholder_idx']) == str(ph_idx):
                            placeholder_info = ph
                            break
            if not placeholder_info:
                continue

            ph_name = placeholder_info['placeholder_name']
            
            if "title" in ph_name.lower():
                display_name = "Title"
                edit_label = "Edit title"
            else:
                display_name = "Content"
                edit_label = "Edit content"

            st.markdown(f"**{display_name}**")

            # If it's a dict for chart/image, we skip editing
            if isinstance(content, dict):
                if "chart_type" in content:
                    st.info("Chart placeholders are not editable via this interface.")
                    continue
                if "image_key" in content:
                    img_key = content["image_key"]
                    if img_key in uploaded_images:
                        st.image(uploaded_images[img_key], 
                                 caption=img_key, 
                                 use_container_width=True)
                        st.info("Image placeholders are not editable via this interface.")
                    continue

                # Otherwise, a dict with text/bullets
                text_val = content.get("text", "")
                bullet_vals = content.get("bullets", [])
            else:
                # Plain string
                text_val = content
                bullet_vals = []

            # Editable text
            edited_text = st.text_area(
                edit_label,
                value=text_val,
                key=f"text_{selected_slide_index}_{ph_idx}",
                height=100
            )

            # Editable bullets
            edited_bullets = []
            if bullet_vals:
                bullets_text = "\n".join(bullet_vals)
                edited_bullets_text = st.text_area(
                    "Edit bullet points (one per line)",
                    value=bullets_text,
                    key=f"bullets_{selected_slide_index}_{ph_idx}"
                )
                edited_bullets = [
                    line.strip() for line in edited_bullets_text.split("\n") if line.strip()
                ]

            # Update placeholders in session state
            if bullet_vals or edited_bullets:
                if isinstance(content, dict):
                    slides[selected_slide_index]['placeholders'][ph_idx]['text'] = edited_text
                    slides[selected_slide_index]['placeholders'][ph_idx]['bullets'] = edited_bullets
                else:
                    slides[selected_slide_index]['placeholders'][ph_idx] = {
                        "text": edited_text,
                        "bullets": edited_bullets
                    }
            else:
                if isinstance(content, dict):
                    slides[selected_slide_index]['placeholders'][ph_idx]['text'] = edited_text
                else:
                    slides[selected_slide_index]['placeholders'][ph_idx] = edited_text

        st.session_state.slides_json['slides'] = slides

        st.markdown("## 3. Download Your Presentation")
        
        # Generate PPT from the updated JSON
        prs = Presentation(EXETER_TEMPLATE_PATH)
        prs = create_slides_from_json(
            prs, 
            st.session_state.slides_json, 
            st.session_state.layout_info, 
            uploaded_images
        )
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        st.download_button(
            "DOWNLOAD",
            data=ppt_buffer.getvalue(),
            file_name="my_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        
        st.markdown("<div style='padding-bottom: 100px'></div>", unsafe_allow_html=True)

if __name__ == "__main__":
    main()
